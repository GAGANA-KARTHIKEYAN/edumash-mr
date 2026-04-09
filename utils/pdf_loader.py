# utils/pdf_loader.py
# Loads PDFs, PPTXs, or plain text curriculum files into LangChain documents

import os
import re


def _Document(page_content, metadata):
    """Version-safe Document factory."""
    try:
        from langchain_core.documents import Document
    except ImportError:
        from langchain.schema import Document
    return Document(page_content=page_content, metadata=metadata)


def _clean_extracted_text(text: str) -> str:
    """Fix common PDF extraction artifacts."""
    # 1. Join hyphenated words broken across lines
    text = re.sub(r"(\w+)-\s*\n\s*(\w+)", r"\1\2", text)
    # 2. Fix multiple spaces/newlines
    text = re.sub(r"\s+", " ", text)
    # 3. Fix common mangled words (e.g. Lpus -> ALUs)
    text = re.sub(r"\bLpus\b", "ALUs", text, flags=re.IGNORECASE)
    # 4. Remove standalone artifacts like "page 1 of 5" or "figure 2"
    text = re.sub(r"Page \d+ of \d+", "", text, flags=re.IGNORECASE)
    return text.strip()


def _is_truly_junk(text: str) -> bool:
    """
    Only block text that is TRULY unreadable:
    - Completely empty
    - Mostly non-printable binary characters (< 50% printable)
    We removed the 'common words' filter because academic/technical/
    non-English PDFs do NOT contain everyday English filler words.
    """
    cleaned = text.strip()
    if not cleaned:
        return True
    # Count printable alphanumeric + punctuation + spaces
    printable = sum(1 for c in cleaned if c.isprintable())
    ratio = printable / len(cleaned)
    if ratio < 0.5:
        print(f"[pdf_loader] ⛔ Blocking junk chunk (printable ratio={ratio:.2f})")
        return True
    # Block very short pages that are just page numbers / headers
    if len(cleaned) < 30:
        return True
    return False


def load_pdf(path: str):
    """Load a PDF file. pdfplumber is primary (better spacing), PyPDFLoader is fallback."""
    print(f"[pdf_loader] Loading PDF: {path}")
    docs = []

    # Primary: pdfplumber - preserves word spacing correctly in academic papers
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                raw_text = page.extract_text(x_tolerance=2, y_tolerance=3) or ""
                text = _clean_extracted_text(raw_text)
                if text.strip():
                    docs.append(_Document(page_content=text, metadata={"source": path, "page": i}))
        print(f"[pdf_loader] pdfplumber: {len(docs)} page(s) from {os.path.basename(path)}")
    except Exception as e:
        print(f"[pdf_loader] pdfplumber failed: {e}")

    # Fallback: PyPDFLoader
    if not docs:
        try:
            from langchain_community.document_loaders import PyPDFLoader
            loader = PyPDFLoader(path)
            docs = loader.load()
            print(f"[pdf_loader] PyPDFLoader fallback: {len(docs)} page(s)")
        except Exception as e:
            print(f"[pdf_loader] PyPDFLoader also failed: {e}")

    return docs



def load_text_fallback(path: str):
    """Load a plain .txt file as a single document."""
    content = ""
    for enc in ("utf-8", "latin-1", "cp1252"):
        try:
            with open(path, "r", encoding=enc, errors="ignore") as f:
                content = f.read()
            break
        except Exception:
            continue
    print(f"[pdf_loader] TXT fallback: {len(content)} chars from {os.path.basename(path)}")
    return [_Document(page_content=content, metadata={"source": path})]


def load_pptx(path: str):
    """Load a PPTX file and return LangChain Document list."""
    print(f"[pdf_loader] Loading PPTX: {path}")
    try:
        import pptx
        prs = pptx.Presentation(path)
        content_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                content_parts.append(f"[Slide {slide_num}] " + " ".join(slide_text))
        full_text = "\n".join(content_parts)
        print(f"[pdf_loader] PPTX: extracted {len(content_parts)} slides, {len(full_text)} chars")
        if not full_text.strip():
            return []
        return [_Document(page_content=full_text, metadata={"source": path})]
    except Exception as e:
        print(f"[pdf_loader] PPTX load failed: {e}")
        return []


def load_all_documents(data_dir: str = "data/"):
    """Load all PDFs, PPTXs, and .txt files from the data/ folder."""
    docs = []
    if not os.path.exists(data_dir):
        print(f"[pdf_loader] data_dir '{data_dir}' does not exist!")
        return docs

    all_files = os.listdir(data_dir)
    print(f"[pdf_loader] Found files in {data_dir}: {all_files}")

    for fname in all_files:
        fpath = os.path.join(data_dir, fname)
        if not os.path.isfile(fpath):
            continue

        fname_lower = fname.lower()
        if fname_lower.endswith(".pdf"):
            raw_docs = load_pdf(fpath)
        elif fname_lower.endswith(".pptx") or fname_lower.endswith(".ppt"):
            raw_docs = load_pptx(fpath)
        elif fname_lower.endswith(".txt"):
            raw_docs = load_text_fallback(fpath)
        else:
            print(f"[pdf_loader] Skipping unsupported file: {fname}")
            continue

        # Only block TRULY unreadable binary garbage
        valid_docs = [d for d in raw_docs if not _is_truly_junk(d.page_content)]
        rejected = len(raw_docs) - len(valid_docs)
        if rejected:
            print(f"[pdf_loader] Rejected {rejected} junk page(s) from {fname}")
        docs.extend(valid_docs)

    print(f"[pdf_loader] ✅ Total readable chunks loaded: {len(docs)}")
    return docs
